#!/usr/bin/env python3
"""
MechAI 统一后端 — 单进程，无需 Docker
SQLite + 内存向量库 + DeepSeek API
"""
import os
import sys
import json
import io
import re
import logging
from typing import List, Optional
from datetime import datetime, timedelta, timezone

import httpx
import numpy as np
from fastapi import FastAPI, Depends, UploadFile, File, Form, Query, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import StreamingResponse
from fastapi.security import HTTPBearer, HTTPAuthorizationCredentials
from pydantic import BaseModel, Field
from sqlalchemy import create_engine, Column, Integer, String, Text, DateTime, Boolean, JSON, ForeignKey, Float
from sqlalchemy.orm import sessionmaker, declarative_base, Session
from jose import jwt, JWTError
from passlib.context import CryptContext

# ===== Logging =====
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger("mechai")

# ===== Config =====
DEEPSEEK_API_KEY = os.getenv("DEEPSEEK_API_KEY", "sk-c7d3bbb1c7a3436890cfec3049fed37b")
DEEPSEEK_BASE_URL = os.getenv("DEEPSEEK_BASE_URL", "https://api.deepseek.com/v1")
SECRET_KEY = os.getenv("SECRET_KEY", "mech-ai-dev-secret-key-change-in-prod")
ALGORITHM = "HS256"
ACCESS_TOKEN_EXPIRE_MINUTES = 1440

# ===== Database (SQLite) =====
DB_PATH = os.path.join(os.path.dirname(os.path.abspath(__file__)), "mechai.db")
engine = create_engine(f"sqlite:///{DB_PATH}", echo=False)
SessionLocal = sessionmaker(autocommit=False, autoflush=False, bind=engine)
Base = declarative_base()

def get_db():
    db = SessionLocal()
    try:
        yield db
    finally:
        db.close()

# ===== ORM Models =====
class User(Base):
    __tablename__ = "users"
    id = Column(Integer, primary_key=True, index=True)
    username = Column(String(50), unique=True, index=True, nullable=False)
    email = Column(String(100), unique=True, index=True, nullable=False)
    hashed_password = Column(String(200), nullable=False)
    full_name = Column(String(100), nullable=True)
    role = Column(String(20), default="user")
    is_active = Column(Boolean, default=True)
    created_at = Column(DateTime, default=datetime.utcnow)

class Document(Base):
    __tablename__ = "documents"
    id = Column(Integer, primary_key=True, index=True)
    title = Column(String(300), nullable=False)
    file_type = Column(String(20), nullable=False)
    file_size = Column(Integer, default=0)
    category = Column(String(100), nullable=True)
    tags = Column(JSON, default=list)
    user_id = Column(Integer, index=True, nullable=False)
    chunk_count = Column(Integer, default=0)
    created_at = Column(DateTime, default=datetime.utcnow)

class DocumentChunk(Base):
    __tablename__ = "document_chunks"
    id = Column(Integer, primary_key=True, index=True)
    document_id = Column(Integer, index=True, nullable=False)
    chunk_index = Column(Integer, nullable=False)
    content = Column(Text, nullable=False)

class Conversation(Base):
    __tablename__ = "conversations"
    id = Column(Integer, primary_key=True, index=True)
    user_id = Column(Integer, index=True, nullable=False)
    title = Column(String(300), default="New Chat")
    created_at = Column(DateTime, default=datetime.utcnow)

class Message(Base):
    __tablename__ = "messages"
    id = Column(Integer, primary_key=True, index=True)
    conversation_id = Column(Integer, index=True, nullable=False)
    role = Column(String(20), nullable=False)
    content = Column(Text, nullable=False)
    created_at = Column(DateTime, default=datetime.utcnow)

# Create tables
Base.metadata.create_all(bind=engine)
logger.info(f"Database ready: {DB_PATH}")

# ===== Security =====
pwd_context = CryptContext(schemes=["bcrypt"], deprecated="auto")
security = HTTPBearer(auto_error=False)

def get_password_hash(pw: str) -> str:
    return pwd_context.hash(pw)

def verify_password(plain: str, hashed: str) -> bool:
    return pwd_context.verify(plain, hashed)

def create_access_token(data: dict) -> str:
    to_encode = data.copy()
    to_encode["exp"] = datetime.now(timezone.utc) + timedelta(minutes=ACCESS_TOKEN_EXPIRE_MINUTES)
    return jwt.encode(to_encode, SECRET_KEY, algorithm=ALGORITHM)

def verify_token(token: str) -> Optional[dict]:
    try:
        return jwt.decode(token, SECRET_KEY, algorithms=[ALGORITHM])
    except JWTError:
        return None

def get_current_user_id(creds: HTTPAuthorizationCredentials = Depends(security)) -> int:
    if not creds:
        raise HTTPException(status_code=401, detail="未登录")
    payload = verify_token(creds.credentials)
    if not payload or not payload.get("sub"):
        raise HTTPException(status_code=401, detail="Token 无效")
    return int(payload["sub"])

# ===== DeepSeek API Client =====
async def deepseek_chat(messages: list, temperature=0.3, max_tokens=2000) -> str:
    async with httpx.AsyncClient(timeout=60.0) as client:
        resp = await client.post(
            f"{DEEPSEEK_BASE_URL}/chat/completions",
            headers={"Authorization": f"Bearer {DEEPSEEK_API_KEY}", "Content-Type": "application/json"},
            json={"model": "deepseek-chat", "messages": messages, "temperature": temperature, "max_tokens": max_tokens},
        )
        resp.raise_for_status()
        return resp.json()["choices"][0]["message"]["content"]

async def deepseek_chat_stream(messages: list, temperature=0.3, max_tokens=2000):
    async with httpx.AsyncClient(timeout=120.0) as client:
        async with client.stream(
            "POST",
            f"{DEEPSEEK_BASE_URL}/chat/completions",
            headers={"Authorization": f"Bearer {DEEPSEEK_API_KEY}", "Content-Type": "application/json"},
            json={"model": "deepseek-chat", "messages": messages, "temperature": temperature, "max_tokens": max_tokens, "stream": True},
        ) as resp:
            resp.raise_for_status()
            async for line in resp.aiter_lines():
                if line.startswith("data: "):
                    data = line[6:].strip()
                    if data == "[DONE]":
                        break
                    try:
                        chunk = json.loads(data)
                        content = chunk["choices"][0].get("delta", {}).get("content", "")
                        if content:
                            yield content
                    except json.JSONDecodeError:
                        continue

async def deepseek_embed(texts: list[str]) -> list[list[float]]:
    """尝试调用 DeepSeek Embedding API，失败则返回 None"""
    try:
        async with httpx.AsyncClient(timeout=15.0) as client:
            resp = await client.post(
                f"{DEEPSEEK_BASE_URL}/embeddings",
                headers={"Authorization": f"Bearer {DEEPSEEK_API_KEY}", "Content-Type": "application/json"},
                json={"model": "deepseek-embedding", "input": texts},
            )
            if resp.status_code == 200:
                data = resp.json()
                sorted_embs = sorted(data["data"], key=lambda x: x["index"])
                return [item["embedding"] for item in sorted_embs]
    except Exception:
        pass
    return None

def cosine_sim(a, b):
    a_np, b_np = np.array(a), np.array(b)
    return float(np.dot(a_np, b_np) / (np.linalg.norm(a_np) * np.linalg.norm(b_np) + 1e-8))

# ===== 本地 TF-IDF 检索 (不依赖外部 Embedding API) =====
import math
from collections import Counter

def _tokenize(text: str) -> list[str]:
    """简单中文分词：按字 + 2-gram"""
    text = re.sub(r'[^\u4e00-\u9fff\w]', ' ', text)
    chars = list(text.replace(' ', ''))
    bigrams = [chars[i]+chars[i+1] for i in range(len(chars)-1)]
    return [c for c in chars if c.strip()] + bigrams

def _tfidf_vector(tokens: list[str], idf: dict[str, float]) -> dict[str, float]:
    tf = Counter(tokens)
    total = len(tokens) or 1
    return {t: (c / total) * idf.get(t, 1.0) for t, c in tf.items()}

def _vec_cosine(a: dict, b: dict) -> float:
    common = set(a) & set(b)
    if not common:
        return 0.0
    dot = sum(a[k] * b[k] for k in common)
    norm_a = math.sqrt(sum(v*v for v in a.values()))
    norm_b = math.sqrt(sum(v*v for v in b.values()))
    return dot / (norm_a * norm_b + 1e-8)

class LocalRetriever:
    """基于 TF-IDF 的本地检索，无需外部 Embedding API"""
    def __init__(self):
        self.docs: list[dict] = []  # {doc_id, chunk_id, content, tokens, vec, user_id, ...}
        self.idf: dict[str, float] = {}

    def _rebuild_idf(self):
        n = len(self.docs)
        if n == 0:
            return
        df = Counter()
        for d in self.docs:
            df.update(set(d["tokens"]))
        self.idf = {t: math.log((n + 1) / (c + 1)) + 1 for t, c in df.items()}
        for d in self.docs:
            d["vec"] = _tfidf_vector(d["tokens"], self.idf)

    def add(self, doc_id: int, chunks: list[str], user_id: int, file_type: str, category: str):
        for i, content in enumerate(chunks):
            tokens = _tokenize(content)
            self.docs.append({
                "doc_id": doc_id, "chunk_id": i, "content": content[:4000],
                "tokens": tokens, "vec": {}, "user_id": user_id,
                "file_type": file_type, "category": category,
            })
        self._rebuild_idf()
        logger.info(f"LocalRetriever: added {len(chunks)} chunks for doc {doc_id}, total={len(self.docs)}")

    def search(self, query: str, top_k=5, user_id=None) -> list[dict]:
        if not self.docs:
            return []
        q_tokens = _tokenize(query)
        q_vec = _tfidf_vector(q_tokens, self.idf)
        results = []
        for d in self.docs:
            if user_id is not None and d["user_id"] != user_id:
                continue
            score = _vec_cosine(q_vec, d["vec"])
            if score >= 0.05:
                results.append({**{k: v for k, v in d.items() if k not in ("tokens", "vec")}, "score": round(score, 4)})
        results.sort(key=lambda x: x["score"], reverse=True)
        return results[:top_k]

    def remove_doc(self, doc_id: int):
        self.docs = [d for d in self.docs if d["doc_id"] != doc_id]
        self._rebuild_idf()

vector_store = LocalRetriever()

# ===== Document Parser =====
def parse_document(file_bytes: bytes, filename: str) -> str:
    ext = os.path.splitext(filename)[1].lower()
    try:
        if ext == ".pdf":
            from PyPDF2 import PdfReader
            reader = PdfReader(io.BytesIO(file_bytes))
            return "\n".join([p.extract_text() or "" for p in reader.pages])
        elif ext in (".docx", ".doc"):
            from docx import Document as DocxDoc
            doc = DocxDoc(io.BytesIO(file_bytes))
            return "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
        elif ext in (".xlsx", ".xls"):
            from openpyxl import load_workbook
            wb = load_workbook(io.BytesIO(file_bytes), data_only=True)
            rows = []
            for sheet in wb.worksheets:
                rows.append(f"=== {sheet.title} ===")
                for row in sheet.iter_rows(values_only=True):
                    line = " | ".join([str(c) if c else "" for c in row])
                    if line.strip(" |"):
                        rows.append(line)
            return "\n".join(rows)
        elif ext in (".pptx", ".ppt"):
            from pptx import Presentation
            prs = Presentation(io.BytesIO(file_bytes))
            texts = []
            for i, slide in enumerate(prs.slides):
                texts.append(f"--- Slide {i+1} ---")
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        texts.append(shape.text_frame.text)
            return "\n".join(texts)
        else:
            return file_bytes.decode("utf-8", errors="ignore")
    except Exception as e:
        logger.warning(f"Parse error for {filename}: {e}")
        return file_bytes.decode("utf-8", errors="ignore")

def chunk_text(text: str, chunk_size=500, overlap=50) -> list[str]:
    text = re.sub(r'\n{3,}', '\n\n', text.strip())
    if len(text) <= chunk_size:
        return [text] if text else []
    chunks, start = [], 0
    while start < len(text):
        end = start + chunk_size
        chunk = text[start:end].strip()
        if chunk:
            chunks.append(chunk)
        start = end - overlap
    return chunks

# ===== FastAPI App =====
app = FastAPI(title="MechAI", version="1.0.0")
app.add_middleware(CORSMiddleware, allow_origins=["*"], allow_credentials=True, allow_methods=["*"], allow_headers=["*"])

SYSTEM_PROMPT = """你是 MechAI，专业的机械工程 AI 助手。你的知识涵盖：
- 机械设计原理（机构学、材料力学、机械零件等）
- 材料科学与选材建议
- 制造工艺（铸造、锻造、焊接、切削加工等）
- 机械故障诊断与分析
- 工程标准与规范（GB、ISO 等）
请用专业但易懂的方式回答，必要时给出公式、参数或步骤。
如果提供了知识库文档内容，请优先参考文档回答。"""

LEARNING_PROMPTS = {
    "课程辅导": """你是 MechAI 学习辅导老师，专注于机械工程核心课程辅导。
你的教学范围：机械原理、机械设计、材料力学、理论力学、工程材料、流体力学、热力学、控制工程等。
教学风格：
- 先给出概念的直觉理解，再讲严格定义
- 用生活中的例子做类比帮助理解
- 适当给出公式推导，但重点讲物理意义
- 用 ✅ ⚠️ ❌ 标记重点和易错点
- 结尾给一个"小结"或"记忆口诀"
如果用户上传了课程文档，请结合文档内容讲解。""",

    "习题生成": """你是 MechAI 智能出题系统，专为机械工程学生生成练习题。
能力：
- 根据指定知识点生成选择题、填空题、计算题、简答题
- 题目难度可选：基础/中等/困难
- 每道题给出详细解答和知识点标注
输出格式：
1. 题目编号 + 难度标签 [基础] [中等] [困难]
2. 题目内容
3. （隐藏答案，等用户回答后再给出）
如果用户要求批改，给出评分和改进建议。
如果用户上传了教材文档，优先从文档中出题。""",

    "考研辅导": """你是 MechAI 考研辅导专家，专注于机械类考研专业课。
覆盖院校：清华、哈工大、西交、华科、浙大、上交等。
覆盖科目：机械原理、机械设计、材料力学、理论力学、控制工程基础。
功能：
- 重点知识梳理（标注高频考点 ⭐）
- 真题解析（说明解题思路和踩分点）
- 复习规划建议
- 易混淆概念辨析
风格：简洁高效，直击考点，不废话。""",

    "错题分析": """你是 MechAI 错题分析师，帮助学生从错误中学习。
工作方式：
1. 分析用户提供的错题和错误答案
2. 诊断错误类型：概念混淆 / 计算失误 / 方法错误 / 知识盲区
3. 给出正确解法，标注关键步骤
4. 推荐 2-3 道类似题巩固
5. 总结这类题的通用解题策略
输出用以下格式：
❌ 错误原因：...
✅ 正确思路：...
📌 关键知识点：...
🔄 巩固练习：...""",
}

ENGINEERING_PROMPTS = {
    "设计建议": """你是 MechAI 机械设计顾问，拥有 20 年机械设计经验。
能力：
- 根据工况（载荷、转速、温度、环境等）推荐结构方案
- 材料选型建议（考虑强度、硬度、耐磨性、成本）
- 安全系数校核（给出计算过程）
- 常见设计陷阱提醒
- 参考标准（GB、ISO、DIN 等）
输出格式：
📊 工况分析 → 🛠️ 推荐方案 → 📐 关键参数 → ⚠️ 注意事项 → 📚 参考标准
如果用户上传了设计文档或图纸说明，请结合实际需求给出建议。""",

    "选型计算": """你是 MechAI 选型计算工程师，精通标准件选型。
覆盖范围：
- 轴承选型（深沟球、角接触、圆锥滚子、调心滚子等）
- 电机选型（步进、伺服、异步、直流等）
- 减速器选型（行星、蜗轮蜗杆、摆线针轮等）
- 气缸/液压缸选型
- 联轴器、皮带、链条传动
- 导轨、丝杠
工作方式：
1. 收集工况参数（转速、扭矩、寿命要求、安装空间等）
2. 计算当量载荷、寿命等
3. 推荐具体型号和品牌
4. 给出安装和维护建议
输出用表格和公式，清晰可执行。""",

    "BOM 分析": """你是 MechAI BOM 分析师，负责物料清单和成本管理。
能力：
- 根据设计方案生成结构化 BOM 清单
- 估算各零部件成本（材料费 + 加工费 + 外购件）
- 推荐供应商和替代方案
- 识别成本优化机会
BOM 输出格式：
| 序号 | 零件名称 | 材料 | 数量 | 单价(估) | 小计 | 备注 |
如果用户上传了 BOM 文件，请分析并给出优化建议。""",

    "工艺路线": """你是 MechAI 工艺工程师，精通机械加工工艺规划。
能力：
- 根据零件特征（形状、材料、精度要求）制定工艺路线
- 工序安排（粗加工→半精加工→精加工→表面处理）
- 工时估算
- 刀具、夹具、量具推荐
- 关键工序质量控制要点
输出格式：
📋 工艺卡片形式，包含：工序号、工序内容、设备、刀具、切削参数、工时
如果用户上传了零件图纸或工艺文件，请针对性分析。""",

    "DFMA": """你是 MechAI DFMA（面向制造和装配的设计）专家。
审查维度：
1. 零件数量最少化 — 能否合并？
2. 装配方向最少化 — 能否单向装配？
3. 防错设计 — 是否有防呆特征？
4. 公差合理性 — 是否过严/过松？
5. 加工可行性 — 是否有难加工特征？
6. 标准化程度 — 能否用标准件替代？
输出格式：
✅ 设计亮点
⚠️ 改进建议（附具体方案）
❌ 必须修改项
📊 DFMA 评分（1-10）
如果用户上传了设计文档或 3D 模型描述，请进行审查。""",

    "FMEA": """你是 MechAI FMEA（失效模式与影响分析）专家。
工作流程：
1. 识别系统/零部件的所有潜在失效模式
2. 分析每种失效的原因和影响
3. 评估严重度(S)、发生度(O)、探测度(D)
4. 计算 RPN = S × O × D
5. 对高 RPN 项给出预防措施
输出格式：
| 失效模式 | 失效原因 | 影响 | S | O | D | RPN | 建议措施 |
标准：参考 AIAG FMEA 手册（第4版）或 IATF 16949。
如果用户上传了相关文档，请结合实际工况分析。""",
}

# ===== Schemas =====
class UserCreate(BaseModel):
    username: str
    password: str
    email: str
    full_name: str = ""

class LoginRequest(BaseModel):
    username: str
    password: str

class ChatRequest(BaseModel):
    message: str
    conversation_id: Optional[int] = None

class LearningRequest(BaseModel):
    topic: str  # 课程辅导 / 习题生成 / 考研辅导 / 错题分析
    message: str
    conversation_id: Optional[int] = None

class EngineeringRequest(BaseModel):
    topic: str  # 设计建议 / 选型计算 / BOM分析 / 工艺路线 / DFMA / FMEA
    message: str
    conversation_id: Optional[int] = None

# ===== API: Health =====
@app.get("/health")
async def health():
    return {"status": "ok", "service": "mechai-unified"}

# ===== API: Users =====
@app.post("/api/users/register")
async def register(data: UserCreate, db: Session = Depends(get_db)):
    existing = db.query(User).filter((User.username == data.username) | (User.email == data.email)).first()
    if existing:
        raise HTTPException(400, "用户名或邮箱已存在")
    user = User(
        username=data.username, email=data.email,
        hashed_password=get_password_hash(data.password),
        full_name=data.full_name,
    )
    db.add(user)
    db.commit()
    db.refresh(user)
    return {"id": user.id, "username": user.username, "email": user.email}

@app.post("/api/users/login")
async def login(data: LoginRequest, db: Session = Depends(get_db)):
    user = db.query(User).filter(User.username == data.username).first()
    if not user or not verify_password(data.password, user.hashed_password):
        raise HTTPException(401, "用户名或密码错误")
    token = create_access_token({"sub": str(user.id), "username": user.username})
    return {"access_token": token, "token_type": "bearer"}

@app.get("/api/users/me")
async def get_me(user_id: int = Depends(get_current_user_id), db: Session = Depends(get_db)):
    user = db.query(User).filter(User.id == user_id).first()
    if not user:
        raise HTTPException(404, "用户不存在")
    return {"id": user.id, "username": user.username, "email": user.email, "full_name": user.full_name}

# ===== API: Documents =====
@app.post("/api/documents/upload")
async def upload_document(
    file: UploadFile = File(...),
    title: str = Form(""),
    category: str = Form(""),
    user_id: int = Depends(get_current_user_id),
    db: Session = Depends(get_db),
):
    if not file.filename:
        raise HTTPException(400, "没有文件")

    file_bytes = await file.read()
    file_size = len(file_bytes)
    text = parse_document(file_bytes, file.filename)
    chunks = chunk_text(text)

    if not chunks:
        raise HTTPException(400, "文档为空或无法解析")

    import os as _os
    file_type = _os.path.splitext(file.filename)[1].lstrip(".")

    doc = Document(
        title=title or file.filename,
        file_type=file_type, file_size=file_size,
        category=category, user_id=user_id,
        chunk_count=len(chunks),
    )
    db.add(doc)
    db.commit()
    db.refresh(doc)

    # Save chunks
    for i, chunk_content in enumerate(chunks):
        db.add(DocumentChunk(document_id=doc.id, chunk_index=i, content=chunk_content))
    db.commit()

    # Store in local retriever
    try:
        vector_store.add(doc.id, chunks, user_id, file_type, category)
    except Exception as e:
        logger.warning(f"Indexing failed: {e}")

    return {"id": doc.id, "title": doc.title, "file_type": doc.file_type, "file_size": doc.file_size, "chunk_count": doc.chunk_count, "created_at": str(doc.created_at)}

@app.get("/api/documents/")
async def list_documents(user_id: int = Depends(get_current_user_id), db: Session = Depends(get_db)):
    docs = db.query(Document).filter(Document.user_id == user_id).order_by(Document.created_at.desc()).all()
    return [{"id": d.id, "title": d.title, "file_type": d.file_type, "file_size": d.file_size, "category": d.category, "chunk_count": d.chunk_count, "created_at": str(d.created_at)} for d in docs]

@app.delete("/api/documents/{doc_id}")
async def delete_document(doc_id: int, user_id: int = Depends(get_current_user_id), db: Session = Depends(get_db)):
    doc = db.query(Document).filter(Document.id == doc_id, Document.user_id == user_id).first()
    if not doc:
        raise HTTPException(404, "文档不存在")
    db.query(DocumentChunk).filter(DocumentChunk.document_id == doc_id).delete()
    db.delete(doc)
    db.commit()
    vector_store.remove_doc(doc_id)
    return {"message": "已删除"}

# ===== API: Chat =====
@app.post("/api/chat/")
async def chat(req: ChatRequest, user_id: int = Depends(get_current_user_id), db: Session = Depends(get_db)):
    """非流式对话"""
    # Get/create conversation
    if req.conversation_id:
        conv = db.query(Conversation).filter(Conversation.id == req.conversation_id).first()
    else:
        conv = Conversation(user_id=user_id, title=req.message[:50])
        db.add(conv)
        db.commit()
        db.refresh(conv)

    # Save user message
    db.add(Message(conversation_id=conv.id, role="user", content=req.message))
    db.commit()

    # History
    msgs = db.query(Message).filter(Message.conversation_id == conv.id).order_by(Message.created_at).all()
    history = [{"role": m.role, "content": m.content} for m in msgs[:-1]]

    # RAG search (local TF-IDF)
    context = None
    try:
        results = vector_store.search(req.message, top_k=3, user_id=user_id)
        if results:
            context = "\n\n---\n\n".join([f"[来源: 文档#{r['doc_id']}]\n{r['content']}" for r in results])
    except Exception as e:
        logger.warning(f"RAG failed: {e}")

    # Build messages
    messages = [{"role": "system", "content": SYSTEM_PROMPT}]
    if context:
        messages.append({"role": "system", "content": f"以下是知识库相关内容：\n\n{context}"})
    messages.extend(history)
    messages.append({"role": "user", "content": req.message})

    reply = await deepseek_chat(messages)

    db.add(Message(conversation_id=conv.id, role="assistant", content=reply))
    db.commit()

    return {"reply": reply, "conversation_id": conv.id, "has_context": context is not None}

@app.post("/api/chat/stream")
async def chat_stream(req: ChatRequest, user_id: int = Depends(get_current_user_id), db: Session = Depends(get_db)):
    """SSE 流式对话"""
    if req.conversation_id:
        conv = db.query(Conversation).filter(Conversation.id == req.conversation_id).first()
    else:
        conv = Conversation(user_id=user_id, title=req.message[:50])
        db.add(conv)
        db.commit()
        db.refresh(conv)

    db.add(Message(conversation_id=conv.id, role="user", content=req.message))
    db.commit()

    msgs = db.query(Message).filter(Message.conversation_id == conv.id).order_by(Message.created_at).all()
    history = [{"role": m.role, "content": m.content} for m in msgs[:-1]]

    context = None
    sources = []
    try:
        results = vector_store.search(req.message, top_k=3, user_id=user_id)
        if results:
            sources = [{"doc_id": r["doc_id"], "score": r["score"], "content": r["content"][:200]} for r in results]
            context = "\n\n---\n\n".join([f"[来源: 文档#{r['doc_id']}]\n{r['content']}" for r in results])
    except Exception as e:
        logger.warning(f"RAG failed: {e}")

    messages = [{"role": "system", "content": SYSTEM_PROMPT}]
    if context:
        messages.append({"role": "system", "content": f"以下是知识库相关内容：\n\n{context}"})
    messages.extend(history)
    messages.append({"role": "user", "content": req.message})

    async def generate():
        if sources:
            yield f"data: {json.dumps({'type': 'sources', 'sources': sources})}\n\n"

        full_reply = []
        async for chunk in deepseek_chat_stream(messages):
            full_reply.append(chunk)
            yield f"data: {json.dumps({'type': 'content', 'content': chunk, 'conversation_id': conv.id})}\n\n"

        db.add(Message(conversation_id=conv.id, role="assistant", content="".join(full_reply)))
        db.commit()
        yield "data: [DONE]\n\n"

    return StreamingResponse(generate(), media_type="text/event-stream")

@app.get("/api/chat/conversations")
async def list_conversations(user_id: int = Depends(get_current_user_id), db: Session = Depends(get_db)):
    convs = db.query(Conversation).filter(Conversation.user_id == user_id).order_by(Conversation.created_at.desc()).all()
    return [{"id": c.id, "title": c.title, "created_at": str(c.created_at)} for c in convs]

@app.get("/api/chat/conversations/{conv_id}")
async def get_conversation(conv_id: int, user_id: int = Depends(get_current_user_id), db: Session = Depends(get_db)):
    conv = db.query(Conversation).filter(Conversation.id == conv_id).first()
    if not conv:
        raise HTTPException(404, "会话不存在")
    msgs = db.query(Message).filter(Message.conversation_id == conv_id).order_by(Message.created_at).all()
    return {
        "id": conv.id, "title": conv.title, "created_at": str(conv.created_at),
        "messages": [{"id": m.id, "role": m.role, "content": m.content, "created_at": str(m.created_at)} for m in msgs],
    }

# ===== API: Learning =====
@app.post("/api/learning/stream")
async def learning_stream(req: LearningRequest, user_id: int = Depends(get_current_user_id), db: Session = Depends(get_db)):
    """学习辅助 — SSE 流式"""
    system_prompt = LEARNING_PROMPTS.get(req.topic, SYSTEM_PROMPT)

    if req.conversation_id:
        conv = db.query(Conversation).filter(Conversation.id == req.conversation_id).first()
    else:
        conv = Conversation(user_id=user_id, title=f"[{req.topic}] {req.message[:40]}")
        db.add(conv)
        db.commit()
        db.refresh(conv)

    db.add(Message(conversation_id=conv.id, role="user", content=req.message))
    db.commit()

    msgs = db.query(Message).filter(Message.conversation_id == conv.id).order_by(Message.created_at).all()
    history = [{"role": m.role, "content": m.content} for m in msgs[:-1]]

    # RAG context
    context = None
    try:
        results = vector_store.search(req.message, top_k=3, user_id=user_id)
        if results:
            context = "\n\n---\n\n".join([f"[来源: 文档#{r['doc_id']}]\n{r['content']}" for r in results])
    except Exception as e:
        logger.warning(f"RAG failed: {e}")

    messages = [{"role": "system", "content": system_prompt}]
    if context:
        messages.append({"role": "system", "content": f"以下是知识库相关内容：\n\n{context}"})
    messages.extend(history)
    messages.append({"role": "user", "content": req.message})

    async def generate():
        full_reply = []
        async for chunk in deepseek_chat_stream(messages):
            full_reply.append(chunk)
            yield f"data: {json.dumps({'type': 'content', 'content': chunk, 'conversation_id': conv.id})}\n\n"

        db.add(Message(conversation_id=conv.id, role="assistant", content="".join(full_reply)))
        db.commit()
        yield "data: [DONE]\n\n"

    return StreamingResponse(generate(), media_type="text/event-stream")

@app.get("/api/learning/exercises")
async def get_exercises(user_id: int = Depends(get_current_user_id), db: Session = Depends(get_db)):
    """获取用户的练习记录"""
    convs = db.query(Conversation).filter(
        Conversation.user_id == user_id,
        Conversation.title.like("[习题生成]%")
    ).order_by(Conversation.created_at.desc()).limit(20).all()
    return [{"id": c.id, "title": c.title, "created_at": str(c.created_at)} for c in convs]

@app.get("/api/learning/wrong-answers")
async def get_wrong_answers(user_id: int = Depends(get_current_user_id), db: Session = Depends(get_db)):
    """获取用户的错题记录"""
    convs = db.query(Conversation).filter(
        Conversation.user_id == user_id,
        Conversation.title.like("[错题分析]%")
    ).order_by(Conversation.created_at.desc()).limit(20).all()
    return [{"id": c.id, "title": c.title, "created_at": str(c.created_at)} for c in convs]

# ===== API: Engineering =====
@app.post("/api/engineering/stream")
async def engineering_stream(req: EngineeringRequest, user_id: int = Depends(get_current_user_id), db: Session = Depends(get_db)):
    """工程辅助 — SSE 流式"""
    system_prompt = ENGINEERING_PROMPTS.get(req.topic, SYSTEM_PROMPT)

    if req.conversation_id:
        conv = db.query(Conversation).filter(Conversation.id == req.conversation_id).first()
    else:
        conv = Conversation(user_id=user_id, title=f"[{req.topic}] {req.message[:40]}")
        db.add(conv)
        db.commit()
        db.refresh(conv)

    db.add(Message(conversation_id=conv.id, role="user", content=req.message))
    db.commit()

    msgs = db.query(Message).filter(Message.conversation_id == conv.id).order_by(Message.created_at).all()
    history = [{"role": m.role, "content": m.content} for m in msgs[:-1]]

    # RAG context
    context = None
    try:
        results = vector_store.search(req.message, top_k=3, user_id=user_id)
        if results:
            context = "\n\n---\n\n".join([f"[来源: 文档#{r['doc_id']}]\n{r['content']}" for r in results])
    except Exception as e:
        logger.warning(f"RAG failed: {e}")

    messages = [{"role": "system", "content": system_prompt}]
    if context:
        messages.append({"role": "system", "content": f"以下是知识库相关内容：\n\n{context}"})
    messages.extend(history)
    messages.append({"role": "user", "content": req.message})

    async def generate():
        full_reply = []
        async for chunk in deepseek_chat_stream(messages):
            full_reply.append(chunk)
            yield f"data: {json.dumps({'type': 'content', 'content': chunk, 'conversation_id': conv.id})}\n\n"

        db.add(Message(conversation_id=conv.id, role="assistant", content="".join(full_reply)))
        db.commit()
        yield "data: [DONE]\n\n"

    return StreamingResponse(generate(), media_type="text/event-stream")

@app.get("/api/engineering/history")
async def get_engineering_history(topic: str = "", user_id: int = Depends(get_current_user_id), db: Session = Depends(get_db)):
    """获取工程辅助历史"""
    query = db.query(Conversation).filter(Conversation.user_id == user_id)
    if topic:
        query = query.filter(Conversation.title.like(f"[{topic}]%"))
    else:
        engineering_topics = list(ENGINEERING_PROMPTS.keys())
        like_conditions = [Conversation.title.like(f"[{t}]%") for t in engineering_topics]
        from sqlalchemy import or_
        query = query.filter(or_(*like_conditions))
    convs = query.order_by(Conversation.created_at.desc()).limit(20).all()
    return [{"id": c.id, "title": c.title, "created_at": str(c.created_at)} for c in convs]

# ===== Main =====
if __name__ == "__main__":
    import uvicorn
    print("=" * 50)
    print("  MechAI 统一后端启动")
    print(f"  DeepSeek API: {DEEPSEEK_BASE_URL}")
    print(f"  Database: {DB_PATH}")
    print("=" * 50)
    uvicorn.run(app, host="0.0.0.0", port=8000)
