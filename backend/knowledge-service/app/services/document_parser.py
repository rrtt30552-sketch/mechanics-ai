import os
import io
import re
from typing import List, Tuple

# Document parsers


def parse_pdf(file_bytes: bytes) -> str:
    from PyPDF2 import PdfReader
    reader = PdfReader(io.BytesIO(file_bytes))
    text_parts = []
    for page in reader.pages:
        text_parts.append(page.extract_text() or "")
    return "\n".join(text_parts)


def parse_docx(file_bytes: bytes) -> str:
    from docx import Document
    doc = Document(io.BytesIO(file_bytes))
    return "\n".join([p.text for p in doc.paragraphs if p.text.strip()])


def parse_xlsx(file_bytes: bytes) -> str:
    from openpyxl import load_workbook
    wb = load_workbook(io.BytesIO(file_bytes), data_only=True)
    rows = []
    for sheet in wb.worksheets:
        rows.append(f"=== Sheet: {sheet.title} ===")
        for row in sheet.iter_rows(values_only=True):
            row_str = " | ".join([str(c) if c is not None else "" for c in row])
            if row_str.strip(" |"):
                rows.append(row_str)
    return "\n".join(rows)


def parse_pptx(file_bytes: bytes) -> str:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(file_bytes))
    texts = []
    for i, slide in enumerate(prs.slides):
        texts.append(f"--- Slide {i+1} ---")
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
    return "\n".join(texts)


PARSERS = {
    ".pdf": parse_pdf,
    ".docx": parse_docx,
    ".doc": parse_docx,
    ".xlsx": parse_xlsx,
    ".xls": parse_xlsx,
    ".pptx": parse_pptx,
    ".ppt": parse_pptx,
}


def parse_document(file_bytes: bytes, filename: str) -> str:
    ext = os.path.splitext(filename)[1].lower()
    parser = PARSERS.get(ext)
    if not parser:
        # Fallback: try to decode as text
        try:
            return file_bytes.decode("utf-8")
        except Exception:
            return file_bytes.decode("gbk", errors="ignore")
    return parser(file_bytes)


def chunk_text(text: str, chunk_size: int = 500, overlap: int = 50) -> List[str]:
    """Split text into overlapping chunks by character count."""
    text = re.sub(r'\n{3,}', '\n\n', text.strip())
    if len(text) <= chunk_size:
        return [text]

    chunks = []
    start = 0
    while start < len(text):
        end = start + chunk_size
        chunk = text[start:end]
        if chunk.strip():
            chunks.append(chunk.strip())
        start = end - overlap
    return chunks
